import streamlit as st
from PIL import Image
import pytesseract
from pptx import Presentation
import tempfile
import fitz  

def ocr_image(img):
    return pytesseract.image_to_string(img)

def ocr_pdf(pdf_file):
    pdf_document = fitz.open(pdf_file)
    text = ""
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text("text")
    return text

def ocr_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

st.markdown(
    """
    <style>
    .title {
        color: #ff4500;
    }
    </style>
    <style>
    .write {
        color: #008080
        }
    .emoji-container {
        position: fixed;
        bottom: 5px;
        left: 50%;
        text-align: center;
        transform: translateX(-50%);
        font-size: 35px;
    }
    </style>
    """,
    unsafe_allow_html=True
)

st.markdown('<h1 class="title">Gorrila Glasses</h1>', unsafe_allow_html=True)
st.markdown('<h3 class="write">Upload the file that you would like converted.  png, jpg, jpeg, PDF and PPTX are supported</h3>', unsafe_allow_html=True)

uploaded_file = st.file_uploader("Choose a file", type=["png", "jpg", "jpeg", "pdf", "pptx"])

if uploaded_file is not None:
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        temp_file.write(uploaded_file.getvalue())
        temp_file_path = temp_file.name

    if st.button('Run OCR'):
        if uploaded_file.type in ["image/png", "image/jpeg"]:
            img = Image.open(uploaded_file)
            st.image(img, caption="Uploaded Image", use_column_width=True)
            text = ocr_image(img)
            st.header("Extracted Text")
            st.text_area(label='', value=text, height=300)

        elif uploaded_file.type == "application/pdf":
            st.write("Uploaded PDF")
            text = ocr_pdf(temp_file_path)
            st.header("Extracted Text")
            st.text_area(label='', value=text, height=900)

        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            st.write("Uploaded PowerPoint")
            text = ocr_pptx(temp_file_path)
            st.header("Extracted Text")
            st.text_area(label='', value=text, height=300)

st.markdown('<div class="emoji-container">🎱</div>', unsafe_allow_html=True)
